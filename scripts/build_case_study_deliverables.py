#!/usr/bin/env python3
"""
Build case study presentation (PPTX) and PDFs:
  - Narrative: docs/case_study/PROMPT_DRIVEN_SAP_CASE_STUDY.md -> PPTX + PDF
  - Repository handout: CASE_STUDY.md (repo root) -> Markdown copy + PDF in exports/

Requires: pip install ".[case-study]"  (python-pptx, fpdf2)

Usage (from repo root):
  python scripts/build_case_study_deliverables.py
  make case-study-artifacts
"""
from __future__ import annotations

import re
import shutil
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
NARRATIVE_CASE_STUDY_MD = REPO_ROOT / "docs" / "case_study" / "PROMPT_DRIVEN_SAP_CASE_STUDY.md"
REPOSITORY_CASE_STUDY_MD = REPO_ROOT / "CASE_STUDY.md"
EXPORT_DIR = REPO_ROOT / "docs" / "case_study" / "exports"
HANDOUT_MD_NAME = "SAP_Repository_Case_Study_Handout.md"
HANDOUT_PDF_NAME = "SAP_Repository_Case_Study_Handout.pdf"


def _ensure_exports() -> None:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)


def _build_pptx() -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:  # pragma: no cover
        raise SystemExit("Install: pip install 'python-pptx' (or pip install -e \".[case-study]\" )") from e

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]  # title + body

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        st = slide.placeholders[1]
        st.text = subtitle

    def add_bullet_slide(title: str, bullets: list[str], notes: str | None = None) -> None:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
        if notes and slide.notes_slide:
            slide.notes_slide.notes_text_frame.text = notes

    # --- Deck content (aligned with the markdown case study) ---
    add_title_slide(
        "Prompt-driven SAP Fiori delivery",
        "MVP, credible sources, official OpenUI5 API, and an LLM-readable registry\n"
        f"Source: {NARRATIVE_CASE_STUDY_MD.name}",
    )

    add_bullet_slide(
        "Problem statement",
        [
            "Design PDFs and Figma do not uniquely map to valid sap.uxap / sap.m XML.",
            "LLMs invent controls, break aggregations, and ignore compact-density behaviour.",
            "Stakeholders need: faster iteration, lower rework cost, and auditable trace (prompt to spec to XML).",
        ],
    )

    add_bullet_slide(
        "Solution overview (three layers)",
        [
            "Fiori & UX: when to use which pattern (object page, worklist, density).",
            "OpenUI5 API: what is legal (aggregations, properties, child types) — authoritative for generation.",
            "Machine index: ComponentSpec, registry.json, recipes — query + validate, do not guess.",
        ],
    )

    add_bullet_slide(
        "MVP approach (this repository)",
        [
            "Vertical slices, each shippable: registry, API, validation, showcase, OPL subscription demo, governance.",
            "If a slice does not run in make/CI, it is not “done” even if the chat output looks good.",
        ],
    )

    add_bullet_slide(
        "Credible source order",
        [
            "1) SAPUI5 API (pinned version)  2) Fiori pattern docs  3) in-repo view + Makefile  4) runtime (DOM)  5) forums / LLM prose last",
            "Forums and raw model answers are hints only — verify in API or browser.",
        ],
    )

    add_bullet_slide(
        "Official SAP guidelines in practice",
        [
            "Controls exist only in the real namespaces (no fictional widgets).",
            "sap_horizon + sapUiSizeCompact in documented bootstraps for demos in this program.",
            "Agent QA in design review: which control ID and aggregation does this XML line come from? (see DESIGN_REVIEW_CHECKLIST.md).",
        ],
    )

    add_bullet_slide(
        "LLM-readable system",
        [
            "JSON Schema + registry for stable IDs; YAML recipes (object page, list-report, etc.).",
            "Retrieval: inject relevant spec slices, not a whole API dump.",
            "Gates: validate-registry, validate-sap-demo, token-audit for static marketing HTML.",
        ],
    )

    add_bullet_slide(
        "Collaboration",
        [
            "UI5 SMEs: aggregations, RGL, OPL blocks behaviour.",
            "Platform: API, tests, reproducible make targets.",
            "Design / product: Fiori intent and MVP scope for credible demos.",
            "Agent authors: prompts aligned with AGENTS.md and Pandya-style constraints.",
        ],
    )

    add_bullet_slide(
        "Step-by-step: prompt to runnable app (high level)",
        [
            "Intent from business prompt → map to pattern + control list",
            "Retrieve ComponentSpec + one recipe (e.g. object page)",
            "Generate view, controller, model; validate; make + browser",
        ],
    )

    add_bullet_slide(
        "Key pain points and mitigations",
        [
            "Invented controls → registry + API as §1",
            "Two OPL blocks, squeezed tables → one VBox, width=100% in blocks",
            "PI + long text in compact → m:Text + short displayValue + formatter",
            "ERR_CONNECTION_REFUSED → run the correct make target, keep server terminal open",
        ],
    )

    add_bullet_slide(
        "Replication (external users)",
        [
            "Clone repo; python -m venv .venv; pip install -e \".[dev]\"",
            "make build-registry; make all (or at least validate targets you touch)",
            "make demo-ui5 / demo-subscription / run-api as documented in README",
        ],
    )

    add_bullet_slide(
        "Deliverables in this project",
        [
            f"Case study: docs/case_study/{NARRATIVE_CASE_STUDY_MD.name}",
            "This PPTX: docs/case_study/exports/SAP_Prompt_Case_Study.pptx",
            "PDF: same folder, SAP_Prompt_Case_Study.pdf (generated alongside)",
            "Google Slides: upload the .pptx to Drive → Open with Google Slides",
        ],
    )

    out = EXPORT_DIR / "SAP_Prompt_Case_Study.pptx"
    prs.save(out)
    return out


def _table_row_to_plain(line: str) -> str | None:
    """If line is a markdown table row, return plain text; None if not a table row; '' to skip (separator)."""
    s = line.strip()
    if not (s.startswith("|") and s.count("|") >= 2):
        return None
    # GFM separator row
    if re.match(r"^\|[\s\-:|]+\|$", s):
        return ""
    parts = [p.strip() for p in s.split("|")]
    # drop empty leading/trailing from split
    parts = [p for p in parts if p != ""]
    if not parts:
        return ""
    out = []
    for p in parts:
        p = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", p)
        p = re.sub(r"\*\*([^*]+)\*\*", r"\1", p)
        p = re.sub(r"`([^`]+)`", r"\1", p)
        out.append(p)
    return "  |  ".join(out)


def _md_to_plain_text(md: str) -> str:
    """Light cleanup for PDF: headings, links, bold; table rows -> plain lines (content preserved)."""
    lines = md.splitlines()
    out: list[str] = []
    for line in lines:
        tr = _table_row_to_plain(line)
        if tr is None:
            out.append(line)
        elif tr == "":
            continue
        else:
            out.append(tr)
    text = "\n".join(out)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    return text


def _register_unicode_font(pdf) -> str:
    """Return font family name to use (DejaVu or Helvetica)."""
    import fpdf

    try:
        font_dir = Path(fpdf.__file__).parent / "font"
        ttf = font_dir / "DejaVuSans.ttf"
        if ttf.is_file():
            pdf.add_font("DejaVu", "", str(ttf))
            pdf.add_font("DejaVu", "B", str(font_dir / "DejaVuSans-Bold.ttf"))
            return "DejaVu"
    except Exception:
        pass
    return "Helvetica"


def _build_pdf_from_md_file(md_path: Path, pdf_path: Path) -> Path:
    try:
        from fpdf import FPDF
    except ImportError as e:  # pragma: no cover
        raise SystemExit("Install: pip install fpdf2 (or pip install -e \".[case-study]\" )") from e

    if not md_path.is_file():
        raise SystemExit(f"Missing {md_path}")

    raw = md_path.read_text(encoding="utf-8")
    plain = _md_to_plain_text(raw)

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    family = _register_unicode_font(pdf)
    pdf.add_page()
    text_w = pdf.w - pdf.l_margin - pdf.r_margin
    if family == "Helvetica":
        # Core fonts: Latin-1 only — normalize common UTF-8 punctuation
        def _lat1(s: str) -> str:
            return (
                s.replace("—", "-")
                .replace("–", "-")
                .replace("’", "'")
                .replace("“", '"')
                .replace("”", '"')
                .replace("…", "...")
            )

        plain = _lat1(plain)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_left_margin(18)
    pdf.set_right_margin(18)

    for line in plain.splitlines():
        line = line.rstrip()
        if not line.strip():
            pdf.ln(4)
            continue
        if line.strip().startswith("---"):
            continue
        is_tableish = "  |  " in line
        if is_tableish:
            pdf.set_font(family, "", 8)
        elif line.isupper() and len(line) < 60 and len(line) > 3 and not line.startswith("http"):
            pdf.set_font(family, "B", 11)
        else:
            pdf.set_font(family, "", 9)
        if family == "Helvetica":
            line = line.encode("latin-1", "replace").decode("latin-1")
        pdf.multi_cell(text_w, 4.2 if is_tableish else 4.5, line)

    pdf.output(pdf_path)
    return pdf_path


def _build_narrative_pdf() -> Path:
    return _build_pdf_from_md_file(
        NARRATIVE_CASE_STUDY_MD,
        EXPORT_DIR / "SAP_Prompt_Case_Study.pdf",
    )


def _build_repository_handout() -> tuple[Path, Path]:
    """Copy root CASE_STUDY.md to exports/ and build matching PDF."""
    if not REPOSITORY_CASE_STUDY_MD.is_file():
        raise SystemExit(f"Missing {REPOSITORY_CASE_STUDY_MD}")

    md_out = EXPORT_DIR / HANDOUT_MD_NAME
    shutil.copy2(REPOSITORY_CASE_STUDY_MD, md_out)
    pdf_out = _build_pdf_from_md_file(REPOSITORY_CASE_STUDY_MD, EXPORT_DIR / HANDOUT_PDF_NAME)
    return md_out, pdf_out


def main() -> None:
    _ensure_exports()
    pptx_path = _build_pptx()
    pdf_path = _build_narrative_pdf()
    handout_md, handout_pdf = _build_repository_handout()
    print(f"Wrote: {pptx_path}")
    print(f"Wrote: {pdf_path}")
    print(f"Wrote: {handout_md}")
    print(f"Wrote: {handout_pdf}")


if __name__ == "__main__":
    main()
