#!/usr/bin/env python3
"""
Build the SAP LLM Playbook deck: clear sections, current → target narrative, presentation-ready
theme (colors, typography, footers) for screen and projector use.

Requires: pip install ".[case-study]"  (python-pptx)

Output: docs/case_study/exports/SAP_LLM_Playbook_Full.pptx
"""
from __future__ import annotations

import struct
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
OUT = REPO / "docs" / "case_study" / "exports" / "SAP_LLM_Playbook_Full.pptx"
# Screenshots: committed PNGs, or replace via `make playbook-images` (localhost capture)
IMAGES = REPO / "docs" / "case_study" / "exports" / "images"

# —— Theme (Fiori-inspired: deep blue, light canvas, legible text)
_C = {
    "dark": RGBColor(0, 60, 113),  # #003C71
    "dark_2": RGBColor(0, 80, 130),
    "canvas": RGBColor(248, 250, 252),
    "text": RGBColor(40, 44, 52),
    "text_muted": RGBColor(100, 108, 118),
    "on_dark": RGBColor(255, 255, 255),
    "on_dark_sub": RGBColor(200, 215, 230),
    "title_accent": RGBColor(0, 60, 113),
    "problem": RGBColor(160, 45, 45),
    "solution": RGBColor(0, 110, 75),
    "footer": RGBColor(120, 128, 138),
    "footer_on_dark": RGBColor(160, 175, 195),
}
FONT = "Calibri"
FONT_TITLE = "Calibri"

# Slide spec: one of
#   {"kind": "title", "title": str, "subtitle": str}
#   {"kind": "section", "label": str, "title": str, "blurb": str}   # e.g. Part 1
#   {"kind": "content", "title": str, "bullets": list[str]}       # level-0 bullets
#   {"kind": "compare", "title", "left": [], "right": [], "left_header": "...", "right_header": "..."}
#   {"kind": "image", "title", "file": "name.png", "caption": str, "note": str optional}


def _slides() -> list[dict[str, Any]]:
    return [
        # —— Opening
        {
            "kind": "title",
            "title": "SAP Fiori + LLM-friendly delivery",
            "subtitle": "From “vibes in chat” to a verified, readable design system\n"
            "Repository: sapui5-llm-ready  |  See DESIGN.md + docs/LLM_HUMAN_PLAYBOOK.md",
        },
        {
            "kind": "content",
            "title": "What this deck covers",
            "bullets": [
                "Part 1 — Current situation: what usually goes wrong with PDF / Figma / raw LLM",
                "Part 2 — How it should be: the target picture (API + JSON + validation)",
                "Part 3 — The bridge: which documents to read, in what order",
                "Part 3b — Skills, path, screenshots  |  3c — Zero→100%  |  3d — :8087 & :8088  |  3e — the **moment** things work (you / agent / repo)",
                "Part 4 — Day to day: prompts, 7 steps, DO / DON’T, errors",
                "Part 5 — Start: clone, make, tools (Claude, Cursor, ChatGPT) + references",
            ],
        },
        # —— Part 1
        {
            "kind": "section",
            "label": "Part 1",
            "title": "Current situation",
            "blurb": "How many teams work today — and why that breaks for reliable SAPUI5 output.",
        },
        {
            "kind": "content",
            "title": "Typical workflow today (common reality)",
            "bullets": [
                "Product / UX: Figma, decks, Confluence = look & intent (good for people).",
                "Developers: read some guidelines; often jump straight to “ask the LLM to build it”.",
                "LLM: free-form answer — not tied to your UI5 version, aggregations, or registry.",
                "Result: first draft *looks* SAP-like; runtime breaks or is impossible to audit.",
            ],
        },
        {
            "kind": "compare",
            "title": "The gap: design handoff vs what code needs",
            "left_header": "What design usually gives",
            "right_header": "What working SAPUI5 code needs",
            "today": [
                "Pixel-perfect frames, component *names* in English",
                "“Use Object Page” in a paragraph (no valid child list)",
                "No link from each screen line to one API fact",
            ],
            "target": [
                "Every control: real `sap.*` class from OpenUI5 API",
                "Exact *aggregations* (e.g. OPL `blocks`, Table `items`) for that version",
                "Optional audit: prompt → retrieved spec line → this XML line",
            ],
        },
        {
            "kind": "content",
            "title": "What actually breaks (symptoms you’ve seen)",
            "bullets": [
                "Invented or wrong control types (`sap.m.Foo` that does not exist).",
                "Wrong or impossible parent/child (illegal aggregation = broken layout).",
                "“Padding” or overlap fixed with CSS when the real issue is structure (e.g. OPL `blocks`, compact + ProgressIndicator).",
                "No answer to: *which official doc justified this line of view.xml?*",
            ],
        },
        # —— Part 2
        {
            "kind": "section",
            "label": "Part 2",
            "title": "How it should be (target state)",
            "blurb": "A clear stack: Fiori intent + API truth + machine index + checks.",
        },
        {
            "kind": "compare",
            "title": "From chaos → to a controlled chain",
            "left_header": "Current (typical) behaviour",
            "right_header": "Target (how we work in this repo)",
            "today": [
                "Chat-only generation",
                "Figma = mistaken for the aggregation spec",
                "Success = “looks like the picture”",
            ],
            "target": [
                "Pipeline: read specs → retrieve JSON slices → generate → `make` + review",
                "Figma = visual + optional `signals.yaml` — API still wins for XML",
                "Success = valid XML + Fiori-credible + traceable to API / registry",
            ],
        },
        {
            "kind": "content",
            "title": "Target stack (3 layers + glue)",
            "bullets": [
                "Human / product: Fiori patterns — *when* Object Page, forms, list-report (design guidelines).",
                "Hard truth: OpenUI5 / SAPUI5 API (ui5.sap.com) for the version in your `index.html`.",
                "Machine index: `ComponentSpec` + `data/registry.json` (schema: `schemas/component_spec.schema.json`).",
                "Glue: YAML recipes, FastAPI search/get/validate, CI (`make all`), design review with agent QA line.",
            ],
        },
        {
            "kind": "content",
            "title": "What “LLM-readable” means (one sentence)",
            "bullets": [
                "The model gets *retrievable* facts (JSON + small recipe), not a vague mood board — so output can be checked against the same rules every time.",
            ],
        },
        # —— Part 3
        {
            "kind": "section",
            "label": "Part 3",
            "title": "The bridge — where to read (before you code)",
            "blurb": "One reading order. No exceptions for production SAPUI5 views.",
        },
        {
            "kind": "content",
            "title": "Read in this order (short)",
            "bullets": [
                "1) OpenUI5 API — version pinned in your app’s `sap-ui-core.js` / `index.html`.",
                "2) Fiori design (pattern choice: object page, density, a11y intent).",
                "3) This repo: `DESIGN.md` (master) → `AGENTS.md` → `SAPUI5-COMPONENTS.md` → registry *slice* for the controls you need.",
            ],
        },
        {
            "kind": "content",
            "title": "Key files in this repository",
            "bullets": [
                "`DESIGN.md` — master map: tokens, Figma vs API, links to all playbooks.",
                "`SAPUI5-COMPONENTS.md` — where API vs `registry.json` vs `recipes/*.yaml`.",
                "`docs/LLM_HUMAN_PLAYBOOK.md` — full narrative for humans + agents.",
                "`docs/PLAYBOOK_INDEX.md` — one-page navigation to every topic file.",
            ],
        },
        {
            "kind": "content",
            "title": "Figma: right role, wrong role",
            "bullets": [
                "Right: layout, copy length, review with stakeholders; optional `data/figma/signals.yaml`.",
                "Wrong: treating frame names or hex as `sap.m` class names or aggregation rules.",
                "If Figma and API disagree → **API** for the XML; document design debt explicitly. See `docs/FIGMA_VS_MACHINE_TRUTH.md`.",
            ],
        },
        {
            "kind": "content",
            "title": "Design “tokens” here — 3 layers (don’t mix them up)",
            "bullets": [
                "Layer 1 — OpenUI5 demos: `sap_horizon` + `sapUiSizeCompact` in `index.html` (see `specs/tokens/closed-sapui5-fiori-set.md`).",
                "Layer 2 — Figma signals: semantic YAML; does **not** compile to valid `blocks` by itself.",
                "Layer 3 — Static marketing HTML in this repo: only variables from `llm-tokens.css` + `make token-audit`.",
            ],
        },
        {
            "kind": "content",
            "title": "Component JSON (what the LLM should retrieve, not invent)",
            "bullets": [
                "Schema: `schemas/component_spec.schema.json` — e.g. `id`, `props`, `slots` (aggregations), `a11y`, `composition`.",
                "Built data: `data/registry.json` — use **slices** in prompts, not the whole file.",
                "Recipes: `recipes/object-page.yaml` etc. — *pattern* templates, must still match API.",
            ],
        },
        # —— Part 3b: Agent skills, path, localhost demos
        {
            "kind": "section",
            "label": "Part 3b",
            "title": "Agent skills, path & local HTML demos",
            "blurb": "What “skills” are, what we ship, the step chain (MCP → JSON → view.xml) — with localhost screenshots.",
        },
        {
            "kind": "content",
            "title": "What is an “agent skill”?",
            "bullets": [
                "A **skill** is a small, versioned package the IDE (e.g. Cursor, Claude Code) can load: usually a **SKILL.md** plus checklists, links, and *when* to use it — so the model follows the *same* rules every session.",
                "Unlike a one-off chat, a skill is **reusable** org knowledge: e.g. “Object Page layout rules”, “read API before XML”, or “7-step validation”.",
                "This repository ships **playbooks in Markdown** (`docs/*.md`, `AGENTS.md`, `DESIGN.md`); some teams also copy an **OPL / UI5 skill** into **`.cursor/skills/`** (see case study: `PROMPT_DRIVEN_SAP_CASE_STUDY.md`, `LLM_HUMAN_PLAYBOOK.md` § tools).",
            ],
        },
        {
            "kind": "content",
            "title": "What is included (skills + repo content)",
            "bullets": [
                "**In-repo:** `AGENTS.md` (agent rules), `DESIGN.md` (master map), `SAPUI5-COMPONENTS.md` (API vs JSON), `data/registry.json` (ComponentSpec slices), `recipes/*.yaml`, FastAPI search/get/validate, `make` gates (`validate-registry`, `validate-sap-demo`, `token-audit`).",
                "**Ecosystem / optional install:** B2B subscription + enterprise showcase **OpenUI5** apps under `examples/*` — run with `make demo-subscription` (8088) and `make demo-showcase` (8087) for a **real browser** check of the same stack you document.",
                "**MCP (optional):** tools that call the *same* HTTP API (search, getComponentSpec, validate) so Claude Desktop can use **tool calls** instead of pasting entire JSON (see Part 9 and `PROMPTING_MCP_AND_STATIC.md`).",
            ],
        },
        {
            "kind": "content",
            "title": "How we get there — end-to-end path (step-by-step)",
            "bullets": [
                "1 **Version** — Pin OpenUI5 / SAPUI5 in `index.html` / `sap-ui-core.js` (all API and XML refer to *that* version).",
                "2 **Pattern** — Choose Fiori pattern (Object Page, list-report, …) from guidelines + a **recipe** YAML when applicable.",
                "3 **Facts (JSON)** — Fetch **ComponentSpec** slices: MCP/HTTP `search` + `getComponentSpec`, *or* attach 2–3 objects from `registry.json` (never the full file in chat).",
                "4 **Structure** — Map slots/aggregations from JSON + API; **Figma** only for copy/layout signals (`signals.yaml`), not invented `sap.m` class names.",
                "5 **SAP XML + JS** — Generate `view.xml` + controller: **only** real `sap.*` control ids, properties, and child aggregations the spec allows.",
                "6 **Run locally** — `make demo-subscription` / `make demo-showcase` → open **http://127.0.0.1:8088/** and **http://localhost:8087** (see next slides).",
                "7 **Validate** — `make validate-*` + 7-step checklist + *which control + aggregation is this line?* on each XML line.",
                "8 **Design gap** — If Figma and API disagree: **API wins** for runtime; log product/design debt explicitly (see `FIGMA_VS_MACHINE_TRUTH.md`).",
            ],
        },
        {
            "kind": "image",
            "title": "Localhost demo — B2B subscription (port 8088)",
            "file": "demo-subscription-8088.png",
            "caption": "http://127.0.0.1:8088/  ·  `make demo-subscription`  ·  `examples/subscription-billing/webapp`",
            "note": "Illustrative / representative frame. Real pixel capture: run demos, then `python3 scripts/playbook_images.py capture` (see Makefile `playbook-images`).",
        },
        {
            "kind": "image",
            "title": "Localhost demo — enterprise LLM showcase (port 8087)",
            "file": "demo-showcase-8087.png",
            "caption": "http://localhost:8087/  ·  `make demo-showcase`  ·  `examples/enterprise-llm-showcase/webapp`",
            "note": "For live UI: `make demo-showcase` (8087). Replace PNG: `playbook_images.py capture` (Chrome + localhost).",
        },
        {
            "kind": "section",
            "label": "Part 3c",
            "title": "Zero → 100% — workflow you can follow",
            "blurb": "Each phase: goal, gap, fix; prompts: wrong vs best; what “HTML-ready” really means. Full text: `docs/WORKFLOW_ZERO_TO_100.md`.",
        },
        {
            "kind": "content",
            "title": "What “100%” means (not a screenshot alone)",
            "bullets": [
                "A **running** OpenUI5 app: **pinned** version in `index.html` — not XML that *looks* SAP from far away but breaks the runtime.",
                "**Reproduce:** clone, `make demo-*`, same URL; **trace:** each important line → API + (optional) ComponentSpec.",
                "Gates: `make validate-*` where your path says so + **7-step** + *which control + aggregation?* for critical XML.",
                "If Fiori/UX and API disagree: **document** the gap; **API wins** for shippable XML (see `FIGMA_VS_MACHINE_TRUTH.md`).",
            ],
        },
        {
            "kind": "content",
            "title": "Six bands: 0% → 100% (do not skip order)",
            "bullets": [
                "0–20%: Env + **pin UI5** + `DESIGN.md` + API for the pattern; **not** Figma as class names.",
                "20–40%: **Recipe** + **registry slice** or MCP/HTTP `search` / `getComponentSpec` — **not** the whole `registry.json` in chat.",
                "40–60%: `view.xml` + controller: **only** real `sap.*` and **legal** parent/child (aggregations).",
                "60–80%: **Browser** — `make demo-showcase` (8087), `demo-subscription` (8088), `demo-ui5` (8085); fix **structure** before `!important` CSS.",
                "80–95%: **`make validate-*`**, 7-step, line-level QA; **not** “looks good” only.",
                "95–100%: Stakeholder-ready + honest **design-debt** notes; **repro** + **trace** = trust.",
            ],
        },
        {
            "kind": "content",
            "title": "Gap → fix (examples you will hit)",
            "bullets": [
                "Hallucinated `sap.m.*` → use **`search` / API**; **never** invent control ids.",
                "Two OPL `blocks` + long text + **compact** → **split** content / one column in `blocks`; watch **ProgressIndicator** + `displayValue` (see playbooks).",
                "Figma name → XML class name → **wrong**; use `signals.yaml` as **soft** input; **API** for real ids.",
                "`ERR_CONNECTION_REFUSED` → wrong: blame UI5. **Right:** start the `make` target; keep terminal open; hard-refresh browser.",
            ],
        },
        {
            "kind": "compare",
            "title": "Prompts: wrong way vs best way (every sprint)",
            "left_header": "Wrong way (avoid)",
            "right_header": "Best way",
            "today": [
                "“Build the Fiori screen from this image.” (no version, no spec)",
                "Paste 50k lines of `registry.json`",
                "“Fix the overlap in CSS” (before aggregations)",
                "“Ship — it works on my machine” (no `make` / 7 steps)",
            ],
            "target": [
                "“UI5 version = …, pattern = …, here are 2 `ComponentSpec` objects; 4-step + line table.”",
                "MCP or HTTP: **small** `search` + `getComponentSpec` for **only** needed ids",
                "**Aggregations** + form density first; then Horizon / theme classes",
                "`make validate-*` + 7-step + **line** QA; then call it done",
            ],
        },
        {
            "kind": "content",
            "title": "Your demo shots ↔ this repository",
            "bullets": [
                "**LLM readiness / top controls** (search bar, Fields + Table tabs, `top_components_seed.yaml` strip): `make demo-showcase` → `http://localhost:8087/`.",
                "**B2B subscription / billing Object Page:** `make demo-subscription` → `http://127.0.0.1:8088/`.",
                "**Create PO / table / shell** (full validation path): `make demo-ui5` → `http://localhost:8085/` + `make validate-sap-demo` for the PO demo.",
                "Full **phase-by-phase** text: `docs/WORKFLOW_ZERO_TO_100.md` — and **:8087 / :8088** deep dive: `docs/DEMOS_LOCALHOST_8087_8088.md`.",
            ],
        },
        {
            "kind": "section",
            "label": "Part 3d",
            "title": "Localhost :8087 & :8088 — how we got to “done”",
            "blurb": "Each demo: purpose, approach, what broke, fix, and the best path to a trustworthy browser outcome.",
        },
        {
            "kind": "content",
            "title": "http://localhost:8087 — Enterprise LLM showcase",
            "bullets": [
                "**Purpose (outcome):** One **Object Page** that proves “**LLM / registry readiness**” in real SAPUI5: search, info strip, **SimpleForm** (ResponsiveGridLayout) + **Table** + actions — **Fiori Horizon** + **Compact**.",
                "**Approach:** Build only from **`data/top_components_seed.yaml`**-aligned `sap.m` + `form` + `sap.uxap`; two **JSONModels** (`view`, `rows`); **no OData**; business copy = ComponentSpec, Figma signals, **SAP demo validation** (table rows).",
                "**Issue → fix:** Unbounded control mix → **hallucination**; *fix* **constrained** set + explicit **OPL** `blocks`. Form **density** *fix* = tuned `labelSpan*`, `columns*`, not random padding. **Outcome:** a **review-safe** “readiness” screenshot tied to the program.",
            ],
        },
        {
            "kind": "content",
            "title": "http://localhost:8088 (or 127.0.0.1) — B2B subscription & billing OPL",
            "bullets": [
                "**Purpose (outcome):** A **B2B** **Object Page** with **section** nav, **Unsubscribe / Re-subscribe** (**dialogs**), **Progress** + **stages** — **in-memory** state machine (no backend) — *prove* dense OPL + real layout pressure.",
                "**Approach:** `JSONModel` + **deep-cloned** state in `Subscription.controller.js`; OPL **sections** with **ObjectPageSubSection** **`blocks`** (Plan, Approval, Services, Balance).",
                "**Issue → fix (real XML):** **OverflowToolbar** + tall **VBox** in header = **clipping** — *fix* use **HBox** + **FlexItemData** + **one row** (`wrap=false`) for contract + **actions** (see `Subscription.view.xml` comments). **ProgressIndicator** in **compact** — *fix* keep **displayValue** short; long copy in **m:Text** (playbook pattern).",
            ],
        },
        {
            "kind": "compare",
            "title": "8087 vs 8088 — two proofs, one best practice",
            "left_header": "8087 (showcase)",
            "right_header": "8088 (subscription)",
            "today": [
                "Narrow: **top components** from seed; **thin** controller",
                "Single **narrative**: “readiness / registry” in one page",
                "Best for: **stakeholder** “what we can render” in 5 min",
            ],
            "target": [
                "Deep: **OPL** depth, **dialogs**, state transitions",
                "Multi-**section** story + **B2B** density + header layout lessons",
                "Best for: **dev** “how to structure **blocks** and avoid **toolbar** traps”",
            ],
        },
        {
            "kind": "content",
            "title": "Best solution — how both reach a successful final stage",
            "bullets": [
                "**Same foundation:** **Pin** UI5; only **API-legal** controls; **ComponentSpec** / registry for **meaning** — not Figma as aggregation truth.",
                "**Layout before CSS:** correct **`blocks`** and **aggregations**; in 8088 especially: **structural** **HBox** fix **before** `!important` on `sap.uxap` (see OPL case material).",
                "**Evidence:** `make demo-showcase` / `make demo-subscription` + browser; add **`make validate-sap-demo`** for the **PO** path (`8085`) when that’s your **ship** story.",
                "**Long form:** `docs/DEMOS_LOCALHOST_8087_8088.md` — print next to your **annotated** screenshots of :8087 and :8088.",
            ],
        },
        {
            "kind": "content",
            "title": "The moment things start to happen — what was “the thing”?",
            "bullets": [
                "**The flip:** work stops being **unbounded chat** and becomes **grounded**: **pinned** UI5 + **retrieved** `ComponentSpec` (or a **small** registry slice) + **recipe** + **`make`** + **browser** — progress = **evidence**, not only screenshots.",
                "**Your action (typical):** you set **non-negotiables** — e.g. only **API-legal** `sap.*` ids, run **`make demo-*` / `build-registry`** before “done”, and **refuse** lines you cannot trace. *That* choice is when the process **stops floating**.",
                "**Assistant’s role:** draft **faster** **inside** those rails; attach **slices**; it does **not** replace the **OpenUI5** doc — **you** still **accept** or **reject** each merge.",
                "**Repo as decision:** `registry` + `make` = **judge** (repeatable *no* to bad XML). **First felt win:** your target **`make`** is **green** **and** **:8087** / **:8088** / **8085** **loads** Fiori-credibly.",
                "Tell **your** story: `docs/TURNING_POINT_MOMENTUM.md` (fill-in: *your* action, *what* the agent did, *what* `make` caught).",
            ],
        },
        # —— Part 4
        {
            "kind": "section",
            "label": "Part 4",
            "title": "Day to day — LLM integration & workflow",
            "blurb": "Three ways to connect tools; then prompt → validate every time.",
        },
        {
            "kind": "content",
            "title": "Three ways to give the model context",
            "bullets": [
                "A) **MCP / tools (best when available)** — e.g. `getComponentSpec`, `search`, `validateUiPlan` (same logic as the HTTP API).",
                "B) **Static docs** — attach `DESIGN.md`, 2–3 `ComponentSpec` objects, one recipe YAML. Good for ChatGPT / Cursor without a server.",
                "C) **Direct files** — hand-pick JSON for one feature; cap size to avoid context noise.",
            ],
        },
        {
            "kind": "content",
            "title": "4-step prompt (use every time)",
            "bullets": [
                "1 **Read** specs (API + DESIGN.md + small registry slice for your controls).",
                "2 **Understand** the pattern (Object Page, form, list) and data shape.",
                "3 **Generate** view + controller — only real control ids and aggregations.",
                "4 **Validate** — `make` targets + 7-step checklist + design review (agent QA on XML lines).",
            ],
        },
        {
            "kind": "content",
            "title": "7-step LLM validation (summary)",
            "bullets": [
                "1 Existence  2 Property names  3 Tokens (per layer)  4 Spec/recipe rules",
                "5 Composition / aggregations  6 Accessibility  7 Patterns + *Which control id + aggregation is this line?*",
                "Details: `docs/VALIDATION_7STEPS_LLM.md`  |  Merge: `docs/DESIGN_REVIEW_CHECKLIST.md`",
            ],
        },
        {
            "kind": "content",
            "title": "5-step workflow (condensed for stand-ups)",
            "bullets": [
                "Read documentation → Select components → Validate choice against API/schema → Generate code → Validate output (make + review).",
            ],
        },
        {
            "kind": "content",
            "title": "If something fails (typical errors)",
            "bullets": [
                "Component not found in API → you hallucinated; search registry + API only.",
                "Property name mismatch → you used HTML/React names; copy from API for that control.",
                "Hardcoded layout hacks → use structure first (OPL, flex), then theme classes.",
                "More: `docs/ERROR_HANDLING_LLM.md`",
            ],
        },
        {
            "kind": "compare",
            "title": "DO (keep) vs DON’T (avoid)",
            "left_header": "DO",
            "right_header": "DON’T",
            "today": [
                "Pin UI5 version; use small registry slices; run `make validate-*`",
                "One agent QA: which control id + aggregation per XML line?",
                "Split long text off ProgressIndicator in compact (m:Text + short %).",
            ],
            "target": [
                "Treat Figma as the aggregation spec",
                "Paste entire `registry.json` into chat",
                "“Fix” OPL with `!important` on `sap.uxap` before fixing `blocks`",
            ],
        },
        {
            "kind": "content",
            "title": "Why this is the “right” outcome (for this program)",
            "bullets": [
                "Reproducible: same repo, same `make`, same tests — not one-off luck in chat.",
                "Auditable: you can point to API + spec lines for Fiori governance.",
                "Teachable: demos (PO, subscription OPL, showcase) + `examples/llm-playbook-comparison/*.html` for React vs Fiori concepts.",
            ],
        },
        # —— Part 5
        {
            "kind": "section",
            "label": "Part 5",
            "title": "Start here — Monday checklist",
            "blurb": "What to run and what to open first.",
        },
        {
            "kind": "content",
            "title": "First hour in the repo",
            "bullets": [
                "`python3 -m venv .venv` → `pip install -e \".[dev]\"`",
                "`make build-registry && make validate-registry`  (or `make all` if you change code).",
                "Read **`CASE_STUDY.md`** (repo root) for the **full program map** — then `DESIGN.md` + [SAPUI5 API](https://ui5.sap.com/#/api) for your version.",
                "Run a demo: e.g. `make demo-ui5` (8085) or `make demo-subscription` (8088) — or `make demo-playbook-site` for the **static HTML** hub (links to all `docs/`) — see `README.md`.",
            ],
        },
        {
            "kind": "content",
            "title": "Bookmark these local URLs (after you start the servers)",
            "bullets": [
                "http://127.0.0.1:8088/  —  B2B subscription & billing Object Page  →  `make demo-subscription` (two terminals: one per demo if both needed)",
                "http://localhost:8087  —  Enterprise LLM showcase (Object Page, seeded top controls)  →  `make demo-showcase`",
                "**Case study HTML** — run `make demo-playbook-site` (opens browser; **127.0.0.1**, first free port from 8089; `scripts/serve_playbook_case_study.py`)  →  **repo root** so `/docs/*.md` links work",
                "If the page will not load: *connection refused* = server not running; start the matching `make` in repo root and keep that terminal open. Hard-refresh the browser after changes.",
            ],
        },
        {
            "kind": "content",
            "title": "Tool cheat-sheet",
            "bullets": [
                "Claude Desktop — MCP to API/tools when you wire it; else attach `DESIGN.md` + small JSON.",
                "Cursor / Claude Code / Windsurf — repo as workspace; `@DESIGN.md` in prompts.",
                "ChatGPT / Copilot — no repo: upload short extracts + link to public SAP API.",
            ],
        },
        {
            "kind": "content",
            "title": "Regenerate this deck (after you edit the script)",
            "bullets": [
                "`pip install -e \".[case-study]\"`  (for python-pptx)",
                "`make playbook-presentation`",
                "Output: `docs/case_study/exports/SAP_LLM_Playbook_Full.pptx`",
            ],
        },
        # —— Part 6+ (from here: continued journey — problem / solution per step)
        {
            "kind": "section",
            "label": "Part 6",
            "title": "Journey: start → each next step (problem / solution)",
            "blurb": "Repeat this rhythm for every feature: goal → what breaks → what you do instead.",
        },
        {
            "kind": "compare",
            "title": "Step 1 — You start (clone & environment)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "No shared truth; everyone asks the model from memory",
                "Dependencies missing; `make` fails mysteriously",
            ],
            "target": [
                "Clone `sapui5-llm-ready`; `pip install -e \".[dev]\"`",
                "Run `make build-registry && make validate-registry` to prove the pipeline",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 2 — You choose a pattern (screen type)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "“Build an object page” with no template → random structure",
            ],
            "target": [
                "Pick a recipe: `recipes/object-page.yaml`, `list-report.yaml`, etc.",
                "Read Fiori pattern docs + OPL / section / `blocks` rules in the playbook",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 3 — You ground the LLM (context)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "Whole API or whole registry in chat → noise + wrong names",
            ],
            "target": [
                "Attach `DESIGN.md` + 2–3 `ComponentSpec` objects + one recipe",
                "Or: MCP / HTTP `search` + `getComponentSpec` for only needed ids",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 4 — You generate (XML + JS)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "Invented `sap.m.*` or React props pasted into `view.xml`",
                "Two OPL `blocks` when you needed one full-width column",
            ],
            "target": [
                "Use only class names from the API; copy property names from the control’s doc page",
                "One `VBox` in `blocks` when stacking; validate with `getCompositionRules`",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 5 — You run the demo (browser)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "`ERR_CONNECTION_REFUSED` — no static server on the port",
                "Stale view: browser cache; “nothing changed”",
            ],
            "target": [
                "Run the right target and leave it open: e.g. `make demo-subscription` → open http://127.0.0.1:8088/",
                "`make demo-showcase` → open http://localhost:8087",
                "Hard refresh after XML/JS edits; ports are defined in `Makefile` (see `README.md`)",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 6 — You validate (gates)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "“Looks fine” in screenshot but illegal child or a11y gap",
            ],
            "target": [
                "`make validate-sap-demo` (and project checks) + 7-step list",
                "Agent QA: *which control id + aggregation is each line?*",
            ],
        },
        {
            "kind": "compare",
            "title": "Step 7 — Next iteration (design vs runtime)",
            "left_header": "Problem",
            "right_header": "Solution",
            "today": [
                "Stakeholder: “not like Figma” but XML is already API-correct",
            ],
            "target": [
                "Record: API-legal; design delta = product decision, not silent hacks",
                "Update Figma notes or `signals.yaml`; do not break aggregations for pixels",
            ],
        },
        {
            "kind": "section",
            "label": "Part 7",
            "title": "What to read (full shelf)",
            "blurb": "Official SAP first, then this repo, then your org’s addendum.",
        },
        {
            "kind": "content",
            "title": "Read order — Layer A (always)",
            "bullets": [
                "OpenUI5 / SAPUI5 API (ui5.sap.com) — your **exact** version from `index.html`",
                "Fiori design guidelines (experience.sap.com) — pattern and density intent",
            ],
        },
        {
            "kind": "content",
            "title": "Read order — Layer B (this repository)",
            "bullets": [
                "`DESIGN.md` → `SAPUI5-COMPONENTS.md` → `docs/LLM_HUMAN_PLAYBOOK.md`",
                "`docs/PLAYBOOK_INDEX.md` — links to validation, Figma, errors, prompting",
                "`AGENTS.md` — rules for any coding / agent session in this code",
            ],
        },
        {
            "kind": "content",
            "title": "Read when you hit a specific issue",
            "bullets": [
                "Layout / OPL: case study + subscription demo + OPL skill if installed",
                "Figma confusion: `docs/FIGMA_VS_MACHINE_TRUTH.md`",
                "Errors from the model: `docs/ERROR_HANDLING_LLM.md`",
                "Org process: `docs/LLM-READABLE-ADDENDUM.md` (2-page hub text)",
            ],
        },
        {
            "kind": "section",
            "label": "Part 8",
            "title": "JSON structure — ComponentSpec & registry",
            "blurb": "What each machine-readable file is for, at a glance.",
        },
        {
            "kind": "content",
            "title": "ComponentSpec: required top-level fields (schema)",
            "bullets": [
                "`schemaVersion`, `id` (e.g. `sap.m.Button`), `name`, `source`, `category`, `status`",
                "`props` (array), `events`, `slots` (= aggregations / allowed children)",
                "`composition` (rules), `tokens` (refs), `a11y`, `examples`, `versioning`",
                "File: `schemas/component_spec.schema.json` — the contract for any JSON you trust",
            ],
        },
        {
            "kind": "content",
            "title": "What the LLM should use for generation (minimum)",
            "bullets": [
                "`id` — must match a real control in the API for that version",
                "`slots` / aggregations — which child types are legal under which parent",
                "`props` — attribute names in XML must appear here (or in API) with correct type",
                "`intentTags` (optional) — e.g. primary-action, dense, for `search` routing",
            ],
        },
        {
            "kind": "content",
            "title": "registry.json — lifecycle in this project",
            "bullets": [
                "Built: `scripts/extract_sapui5.py` from API docs (fixtures or live)",
                "Optional merge: `apply_figma_signals.py` for semantic metadata — **not** new aggregations from pixels",
                "Validate: `make validate-registry` against the schema; CI blocks broken contracts",
            ],
        },
        {
            "kind": "section",
            "label": "Part 9",
            "title": "MCP (Model Context Protocol) — more detail",
            "blurb": "Tools, not prose, for the same HTTP API the repo already ships.",
        },
        {
            "kind": "content",
            "title": "What MCP is (in one slide)",
            "bullets": [
                "A standard way for a client (e.g. Claude Desktop) to call **tools** with structured input/output",
                "Here: the **interesting** tools mirror `search`, `getComponentSpec`, `validateUiPlan`",
                "Goal: the model **must** use retrieved facts, not “remembered” control names",
            ],
        },
        {
            "kind": "content",
            "title": "Suggested map: API capability → “tool” idea",
            "bullets": [
                "`GET /components/{id}`  →  tool: get_component_spec(id)",
                "`POST /search`  →  tool: search_components(intent, filters)",
                "Validation route  →  tool: validate_ui_plan(json_plan)",
                "Implementation: wrap FastAPI or run MCP server in front of the same `sapui5_llm_ready` code",
            ],
        },
        {
            "kind": "content",
            "title": "Why MCP / tools beat a giant prompt (for SAP)",
            "bullets": [
                "UI5 has hundreds of controls — you inject **only** what the step needs",
                "Returns are **validatable** (JSON) — easier to chain and to test than free text",
                "Reduces **hallucinated** `sap.m.*` ids because the model sees real registry rows",
            ],
        },
        {
            "kind": "content",
            "title": "No MCP server? (fallbacks that still work)",
            "bullets": [
                "Run `make run-api` and call HTTP from a script, or from Cursor terminal",
                "Attach **small** static JSON from `data/registry.json` (curated copy-paste, not 50 MB)",
                "Use the 4-step prompt in `docs/PROMPTING_MCP_AND_STATIC.md`",
            ],
        },
        {
            "kind": "section",
            "label": "Part 10",
            "title": "More to include in your org pack",
            "blurb": "Governance, education assets, and where to go next.",
        },
        {
            "kind": "content",
            "title": "Governance, CI, and quality (details)",
            "bullets": [
                "`GOVERNANCE.md` — what must stay green on merge",
                "`make all` — lint, registry, tests, token-audit, SAP demo validation (see Makefile)",
                "Design review: `docs/DESIGN_REVIEW_CHECKLIST.md`",
            ],
        },
        {
            "kind": "content",
            "title": "Education assets in the repo (details)",
            "bullets": [
                "`docs/case_study/PROMPT_DRIVEN_SAP_CASE_STUDY.md` — long narrative",
                "`examples/llm-playbook-comparison/fiori-concepts-vs-react.html` — same fields, two stacks",
                "Live OpenUI5 in browser: http://127.0.0.1:8088/ (`make demo-subscription`) and http://localhost:8087 (`make demo-showcase`)",
            ],
        },
        {
            "kind": "content",
            "title": "This deck: how it grows (for maintainers)",
            "bullets": [
                "All slides are data in `scripts/build_playbook_presentation.py` → function `_slides()`",
                "Add more dicts: `content`, `compare`, `section`, `title` — then `make playbook-presentation`",
                "Outline for editors: `docs/case_study/PLAYBOOK_DECK_OUTLINE.md`",
            ],
        },
        {
            "kind": "title",
            "title": "Thank you",
            "subtitle": "Full detail: docs/LLM_HUMAN_PLAYBOOK.md\n"
            "Questions: API version, registry slice, one recipe per feature, then validate.\n"
            "Part 3b–e: skills, :8087/:8088, zero→100%, **turning point** (TURNING_POINT_MOMENTUM.md). Parts 6–10: journey, JSON, MCP, governance.",
        },
    ]


def _set_bg(s, rgb: RGBColor) -> None:
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def _p_font(p, *, size: Pt, bold: bool = False, color: RGBColor | None = None, name: str = FONT) -> None:
    p.font.size = size
    p.font.bold = bold
    p.font.name = name
    if color is not None:
        p.font.color.rgb = color
    for r in p.runs:
        r.font.size = size
        r.font.bold = bold
        r.font.name = name
        if color is not None:
            r.font.color.rgb = color


def _footer(s, on_dark: bool) -> None:
    rgb = _C["footer_on_dark"] if on_dark else _C["footer"]
    box = s.shapes.add_textbox(
        Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.45)
    )
    tfm = box.text_frame
    tfm.text = "sapui5-llm-ready  ·  SAP Fiori + LLM-readable delivery  ·  Screen / projector"
    p = tfm.paragraphs[0]
    _p_font(p, size=Pt(9), color=rgb, bold=False)
    p.alignment = PP_ALIGN.CENTER


def _add_title(prs: Presentation, title: str, subtitle: str) -> None:
    lo = prs.slide_layouts[0]
    s = prs.slides.add_slide(lo)
    _set_bg(s, _C["dark"])
    s.shapes.title.text = title
    t = s.shapes.title.text_frame
    t.word_wrap = True
    t.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in t.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        _p_font(p, size=Pt(40), bold=True, color=_C["on_dark"], name=FONT_TITLE)
    if len(s.placeholders) > 1:
        ph1 = s.placeholders[1]
        ph1.text = subtitle
        t2 = ph1.text_frame
        t2.word_wrap = True
        for i, p in enumerate(t2.paragraphs):
            p.alignment = PP_ALIGN.CENTER
            _p_font(
                p,
                size=Pt(20) if i == 0 else Pt(18),
                color=_C["on_dark_sub"],
            )
    _footer(s, on_dark=True)


def _add_section(prs: Presentation, label: str, title: str, blurb: str) -> None:
    lo = prs.slide_layouts[0]
    s = prs.slides.add_slide(lo)
    _set_bg(s, _C["dark_2"])
    s.shapes.title.text = f"{label}\n{title}"
    t = s.shapes.title.text_frame
    t.word_wrap = True
    t.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in t.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        _p_font(p, size=Pt(32), bold=True, color=_C["on_dark"], name=FONT_TITLE)
    s.placeholders[1].text = blurb
    sub = s.placeholders[1].text_frame
    sub.word_wrap = True
    for p in sub.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        _p_font(p, size=Pt(20), color=_C["on_dark_sub"])
    _footer(s, on_dark=True)


def _add_content(prs: Presentation, title: str, bullets: list[str]) -> None:
    lo = prs.slide_layouts[1]
    s = prs.slides.add_slide(lo)
    _set_bg(s, _C["canvas"])
    s.shapes.title.text = title
    pti = s.shapes.title.text_frame.paragraphs[0]
    _p_font(
        pti, size=Pt(28), bold=True, color=_C["title_accent"], name=FONT_TITLE
    )
    tf = s.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.level = 0
    _p_font(p, size=Pt(18), color=_C["text"])
    p.space_after = Pt(6)
    for line in bullets[1:]:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.level = 0
        _p_font(p2, size=Pt(18), color=_C["text"])
        p2.space_after = Pt(6)
    tf.word_wrap = True
    for para in tf.paragraphs:
        if para.text.strip():
            try:
                para.line_spacing = 1.2
            except Exception:  # pragma: no cover
                pass
    _footer(s, on_dark=False)


def _add_compare(
    prs: Presentation,
    title: str,
    today: list[str],
    target: list[str],
    left_header: str = "Current situation (today)",
    right_header: str = "How it should be (target)",
) -> None:
    lo = prs.slide_layouts[1]
    s = prs.slides.add_slide(lo)
    _set_bg(s, _C["canvas"])
    s.shapes.title.text = title
    _p_font(
        s.shapes.title.text_frame.paragraphs[0],
        size=Pt(26),
        bold=True,
        color=_C["title_accent"],
        name=FONT_TITLE,
    )
    tf = s.shapes.placeholders[1].text_frame
    tf.text = left_header
    header_map_left = {
        "Problem": _C["problem"],
        "DO": RGBColor(0, 100, 60),
    }
    h1 = header_map_left.get(left_header, _C["title_accent"])
    _p_font(tf.paragraphs[0], size=Pt(17), bold=True, color=h1, name=FONT)
    gap = 3 if left_header in ("Problem", "DO") else 4
    for line in today:
        p2 = tf.add_paragraph()
        p2.text = f"• {line}"
        p2.level = 0
        _p_font(p2, size=Pt(16), color=_C["text"])
        p2.space_after = Pt(gap)
    p_sp = tf.add_paragraph()
    p_sp.text = ""
    p_sp.space_after = Pt(6)
    p_mid = tf.add_paragraph()
    p_mid.text = right_header
    h2 = {"Solution": _C["solution"], "DON’T": RGBColor(150, 50, 50)}.get(
        right_header, _C["title_accent"]
    )
    _p_font(p_mid, size=Pt(17), bold=True, color=h2, name=FONT)
    p_mid.space_after = Pt(4)
    for line in target:
        p3 = tf.add_paragraph()
        p3.text = f"• {line}"
        p3.level = 0
        _p_font(p3, size=Pt(16), color=_C["text"])
    tf.word_wrap = True
    _footer(s, on_dark=False)


def _png_size(path: Path) -> tuple[int, int] | None:
    """Read width/height from a PNG (first IHDR) without third-party deps."""
    try:
        with path.open("rb") as f:
            sig = f.read(8)
            if len(sig) < 8 or sig[1:4] != b"PNG":
                return None
            plen = int.from_bytes(f.read(4), "big")
            typ = f.read(4)
            if typ != b"IHDR" or plen < 8:
                return None
            data = f.read(8)
            if len(data) < 8:
                return None
            w, h = struct.unpack(">II", data)
            return w, h
    except (OSError, struct.error, ValueError):
        return None


def _add_image(
    prs: Presentation,
    title: str,
    file_name: str,
    caption: str,
    note: str | None = None,
) -> None:
    path = (IMAGES / file_name).resolve()
    lo = prs.slide_layouts[5]  # Title Only
    s = prs.slides.add_slide(lo)
    _set_bg(s, _C["canvas"])
    s.shapes.title.text = title
    _p_font(
        s.shapes.title.text_frame.paragraphs[0],
        size=Pt(22),
        bold=True,
        color=_C["title_accent"],
        name=FONT_TITLE,
    )

    if not path.is_file():
        err = s.shapes.add_textbox(
            Inches(0.55), Inches(2.0), Inches(12.2), Inches(2.0)
        )
        err.text_frame.text = (
            f"Image not found: {path}\n"
            f"Add the PNG to docs/case_study/exports/images/ or run: "
            f"make playbook-images (see scripts/playbook_images.py)"
        )
        for p in err.text_frame.paragraphs:
            _p_font(p, size=Pt(14), color=_C["problem"], bold=True)
    else:
        dim = _png_size(path) or (1400, 780)
        w_px, h_px = dim
        max_w, max_h = 12.1, 5.05
        ar = w_px / h_px
        box_ar = max_w / max_h
        if ar > box_ar:
            pic_w, pic_h = max_w, max_w / ar
        else:
            pic_h = max_h
            pic_w = max_h * ar
        left = (13.333 - pic_w) / 2.0
        top = 0.95
        s.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(pic_w), Inches(pic_h)
        )

    cap = s.shapes.add_textbox(Inches(0.45), Inches(6.28), Inches(12.4), Inches(0.55))
    cap.text_frame.text = caption
    for p in cap.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        _p_font(p, size=Pt(10), color=_C["text_muted"], bold=False)
    if note:
        nbox = s.shapes.add_textbox(Inches(0.45), Inches(6.78), Inches(12.4), Inches(0.35))
        nbox.text_frame.text = note
        for p in nbox.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            _p_font(p, size=Pt(8), color=_C["footer"], bold=False)
    _footer(s, on_dark=False)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in _slides():
        k = spec["kind"]
        if k == "title":
            _add_title(prs, spec["title"], spec["subtitle"])
        elif k == "section":
            _add_section(prs, spec["label"], spec["title"], spec["blurb"])
        elif k == "content":
            _add_content(prs, spec["title"], spec["bullets"])
        elif k == "compare":
            _add_compare(
                prs,
                spec["title"],
                spec["today"],
                spec["target"],
                left_header=spec.get("left_header", "Current situation (today)"),
                right_header=spec.get("right_header", "How it should be (target)"),
            )
        elif k == "image":
            _add_image(
                prs,
                spec["title"],
                spec["file"],
                spec["caption"],
                note=spec.get("note"),
            )
        else:  # pragma: no cover
            raise ValueError(f"Unknown kind {k}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    n = len(prs.slides)
    print(f"Wrote {OUT} ({n} slides)")


if __name__ == "__main__":
    main()
